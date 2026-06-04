import asyncio
import abc
import typing
import os
from pptx import Presentation
from pptx.shapes.base import BaseShape
from pypdf import PdfReader
from pptx.shapes.base import BaseShape
from pptx.shapes.autoshape import Shape
from loguru import logger
import docx
import io
import re

class Analyzer(abc.ABC):
    @staticmethod
    async def analyze(*args, **kwargs) -> typing.Any:
        pass



#=============ВЕРНИСЬ И ДОДЕЛАЙ СИСТЕМУ ВОЗВРАТА В АНАЛИЗАТОРЕ.
@typing.final
class PresentAnalyzer(Analyzer):
    @staticmethod
    async def analyze(presen, words_list: list) -> typing.Any:
        try:
            slide_cout = 0
            extracted_text = ""

            if presen.endswith('.pptx'):
                    prs = Presentation(presen)
                    slide_count = len(prs.slides)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                real_shape = shape
                                text = real_shape.text_frame.text#type: ignore
                                if text:
                                    extracted_text += " " + text.lower()


            elif presen.endswith('.pdf'):
                reader = PdfReader(presen)
                slide_count = len(reader.pages)
                for page in reader.pages:
                    text = page.extract_text()
                    if text:
                        extracted_text += " " + text.lower()
            else:
                logger.info("❌ Неподдерживаемый формат файла!")
                return

            validation_results = {}
            missing_sections = []


            for keyword in words_list:
                is_found = keyword in extracted_text
                validation_results[keyword] = "✅ НАЙДЕНО" if is_found else "❌ ОТСУТСТВУЕТ"
                if not is_found:
                    missing_sections.append(keyword)
                    
        except Exception as e:
            pass


#def analyze_presentation(file_path: str):
#    # 4. Валидация ключевых слов
#    validation_results = {}
#    missing_sections = []
#    
#    for keyword in REQUIRED_KEYWORDS:
#        is_found = keyword in extracted_text
#
#        validation_results[keyword] = "✅ НАЙДЕНО" if is_found else "❌ ОТСУТСТВУЕТ"
#        if not is_found:
#            missing_sections.append(keyword)
#
#    # 5. Вывод результатов в консоль
#    print(f"Всего слайдов/страниц: {slide_count}")
#    print("\nРезультаты проверки структуры:")
#    for criteria, status in validation_results.items():
#        print(f"  - {criteria.capitalize()}: {status}")
#    
#    print("\nВердикт:")
#    if not missing_sections:
#        print("🎉 Успех! Презентация соответствует всем требованиям ТЗ.")
#    else:
#        print(f"⚠️ Внимание! Не хватает следующих разделов: {', '.join(missing_sections)}")
#
# Запуск проверки для нашего файла
#if __name__ == "__main__":
#    # Укажи имя своего файла здесь
#    analyze_presentation("test (1).pdf")
#




@typing.final
class DocumentationAnaluzer(Analyzer): 
    @staticmethod
    async def analyze(words_list: list, min_char: int, filename: str,  file_bytes:bytes) -> typing.Any:
        filename_lower = filename.lower()
        text_content = ""
        image_count = 0

        # 2. Извлечение текста и подсчет изображений в зависимости от формата
        if filename_lower.endswith('.docx'):
            doc = docx.Document(io.BytesIO(file_bytes))
            # Собираем текст из всех параграфов
            text_content = "\n".join([p.text for p in doc.paragraphs])
            # Считаем картинки, встроенные в документ
            image_count = len(doc.inline_shapes)

        elif filename_lower.endswith('.pdf'):
            reader = PdfReader(io.BytesIO(file_bytes))
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    text_content += "\n" + text
                # Пытаемся извлечь количество картинок на странице через встроенные ресурсы pypdf
                image_count += len(page.images)

        elif filename_lower.endswith('.md'):
            # Декодируем байты текстового файла Markdown в строку
            text_content = file_bytes.decode('utf-8', errors='ignore')
            # В Markdown изображения задаются конструкциями ![текст](ссылка) или тегами <img>
            md_images = re.findall(r'!\[.*?\]\(.*?\)', text_content)
            html_images = re.findall(r'<img', text_content)
            image_count = len(md_images) + len(html_images)
            
        else:
            return {"error": f"Формат файла {filename} не поддерживается."}

        # Приводим весь текст к нижнему регистру для поиска
        text_lower = text_content.lower()
        total_chars = len(text_content.strip())

        # 3. Проверка наличия обязательных разделов
        found_sections = {}
        missing_sections = []

        for section_id, keywords in words_list.items():
            is_found = any(keyword in text_lower for keyword in keywords)
            found_sections[section_id] = "✅ Найдено" if is_found else "❌ Отсутствует"
            if not is_found:
                missing_sections.append(section_id)

        # 4. Базовая оценка качества
        has_enough_volume = total_chars >= MIN_CHARACTER_COUNT
        has_diagrams = image_count > 0
        
        # Ищем ссылки на внешние ресурсы/документы (http/https)
        has_links = len(re.findall(r'https?://[^\s]+', text_content)) > 0

        # Итоговый вердикт качества
        quality_score = "Отличная документация"
        if missing_sections:
            quality_score = "Критическая ошибка: отсутствуют обязательные разделы"
        elif not has_enough_volume:
            quality_score = "Низкое качество: слишком маленький объем текста"
        elif not has_diagrams:
            quality_score = "Среднее качество: текст не подкреплен схемами или изображениями"

        return {
            "filename": filename,
            "total_characters": total_chars,
            "detected_images": image_count,
            "has_external_links": has_links,
            "sections_status": found_sections,
            "missing_sections": missing_sections,
            "quality_score": quality_score,
            "is_valid": len(missing_sections) == 0 and has_enough_volume




# 1. Настройки валидации (Критерии организаторов)
# Ищем разделы по ключевым корням слов, чтобы учесть разные падежи и синонимы
REQUIRED_SECTIONS = {
    "описание_работы": ["описание работы", "архитектура системы", "общие сведения", "описание системы"],
    "развертывание": ["развёртывание", "развертывание", "инструкция по установке", "установка", "deployment", "docker", "compose", "venv"],
    "эксплуатация": ["эксплуатация","эксплуатации", "инструкция по эксплуатации", "руководство пользователя", "использование"]
}

MIN_CHARACTER_COUNT = 1000  # Минимальный общий объём текста для качественной документации

def analyze_documentation(file_bytes: bytes, filename: str) -> dict:
    """
    Анализирует бинарные данные документа. 
    Функция универсальна: принимает байты из файла, БД или сети.
    """
    filename_lower = filename.lower()
    text_content = ""
    image_count = 0

    # 2. Извлечение текста и подсчет изображений в зависимости от формата
    if filename_lower.endswith('.docx'):
        doc = docx.Document(io.BytesIO(file_bytes))
        # Собираем текст из всех параграфов
        text_content = "\n".join([p.text for p in doc.paragraphs])
        # Считаем картинки, встроенные в документ
        image_count = len(doc.inline_shapes)

    elif filename_lower.endswith('.pdf'):
        reader = PdfReader(io.BytesIO(file_bytes))
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_content += "\n" + text
            # Пытаемся извлечь количество картинок на странице через встроенные ресурсы pypdf
            image_count += len(page.images)

    elif filename_lower.endswith('.md'):
        # Декодируем байты текстового файла Markdown в строку
        text_content = file_bytes.decode('utf-8', errors='ignore')
        # В Markdown изображения задаются конструкциями ![текст](ссылка) или тегами <img>
        md_images = re.findall(r'!\[.*?\]\(.*?\)', text_content)
        html_images = re.findall(r'<img', text_content)
        image_count = len(md_images) + len(html_images)
        
    else:
        return {"error": f"Формат файла {filename} не поддерживается."}

    # Приводим весь текст к нижнему регистру для поиска
    text_lower = text_content.lower()
    total_chars = len(text_content.strip())

    # 3. Проверка наличия обязательных разделов
    found_sections = {}
    missing_sections = []

    for section_id, keywords in REQUIRED_SECTIONS.items():
        is_found = any(keyword in text_lower for keyword in keywords)
        found_sections[section_id] = "✅ Найдено" if is_found else "❌ Отсутствует"
        if not is_found:
            missing_sections.append(section_id)

    # 4. Базовая оценка качества
    has_enough_volume = total_chars >= MIN_CHARACTER_COUNT
    has_diagrams = image_count > 0
    
    # Ищем ссылки на внешние ресурсы/документы (http/https)
    has_links = len(re.findall(r'https?://[^\s]+', text_content)) > 0

    # Итоговый вердикт качества
    quality_score = "Отличная документация"
    if missing_sections:
        quality_score = "Критическая ошибка: отсутствуют обязательные разделы"
    elif not has_enough_volume:
        quality_score = "Низкое качество: слишком маленький объем текста"
    elif not has_diagrams:
        quality_score = "Среднее качество: текст не подкреплен схемами или изображениями"

    return {
        "filename": filename,
        "total_characters": total_chars,
        "detected_images": image_count,
        "has_external_links": has_links,
        "sections_status": found_sections,
        "missing_sections": missing_sections,
        "quality_score": quality_score,
        "is_valid": len(missing_sections) == 0 and has_enough_volume
    }

if __name__ == "__main__":
    target_file = "readme_perfect.md"
    if os.path.exists(target_file):
        with open(target_file, "rb") as f:
            bytes_data = f.read()
        fake_filename = target_file
    result = analyze_documentation(bytes_data, target_file)
    
    import json
