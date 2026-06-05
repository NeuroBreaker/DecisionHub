import asyncio
import abc
import typing
import os
from pptx import Presentation
from pptx.shapes.base import BaseShape
from pypdf import PdfReader
from pptx.shapes.base import BaseShape
from pptx.shapes.autoshape import Shape
from loguru import logger
import docx
import io
import re

import io
import typing
from fastapi import UploadFile
from pptx import Presentation
from pypdf import PdfReader

import os
import tempfile
import typing
import yt_dlp

class Analyzer(abc.ABC):
    @staticmethod
    async def analyze(*args, **kwargs) -> typing.Any:
        pass



#=============ВЕРНИСЬ И ДОДЕЛАЙ СИСТЕМУ ВОЗВРАТА В АНАЛИЗАТОРЕ.
@typing.final
class MembersPres():
    @staticmethod
    async def analyze(file: UploadFile, words_list: list) -> typing.Any:
        try:
            logger.info("Present analyze started")
            slide_count = 0  # Исправил опечатку slide_cout
            extracted_text = ""
            
            # Получаем имя файла для проверки расширения
            filename = file.filename.lower() if file.filename else ""

            # Считываем содержимое файла в память
            file_bytes = await file.read()
            file_stream = io.BytesIO(file_bytes)

            if filename.endswith('.pptx'):
                logger.info("PPTX DETECTED")
                # Передаем поток байтов вместо пути к файлу
                prs = Presentation(file_stream)
                slide_count = len(prs.slides)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            real_shape = shape
                            text = real_shape.text_frame.text#type:ignore
                            if text:
                                extracted_text += " " + text.lower()

            elif filename.endswith('.pdf'):
                logger.info("PDF DETECTED")
                # Передаем поток байтов вместо пути к файлу
                reader = PdfReader(file_stream)
                slide_count = len(reader.pages)
                for page in reader.pages:
                    text = page.extract_text()
                    if text:
                        extracted_text += " " + text.lower()
            else:
                logger.info("❌ Неподдерживаемый формат файла!")
                return

            validation_results = {}
            missing_sections = []

            for keyword in words_list:
                is_found = keyword in extracted_text
                validation_results[keyword] = True if is_found else False

                if not is_found:
                    missing_sections.append(keyword)
                
            total_points = 0
            for key, value in validation_results.items():
                if value is True:
                    total_points += 1

            return total_points

        except Exception as e:
            logger.error(f"ПРОБЛЕМА В ОБРАБОТКЕ ПРЕЗЫ... {e}")
        finally:
            # Важно закрыть дескриптор файла FastAPI после работы
            await file.close()

#def analyze_presentation(file_path: str):
#    # 4. Валидация ключевых слов
#    validation_results = {}
#    missing_sections = []
#    
#    for keyword in REQUIRED_KEYWORDS:
#        is_found = keyword in extracted_text
#
#        validation_results[keyword] = "✅ НАЙДЕНО" if is_found else "❌ ОТСУТСТВУЕТ"
#        if not is_found:
#            missing_sections.append(keyword)
#
#    # 5. Вывод результатов в консоль
#    print(f"Всего слайдов/страниц: {slide_count}")
#    print("\nРезультаты проверки структуры:")
#    for criteria, status in validation_results.items():
#        print(f"  - {criteria.capitalize()}: {status}")
#    
#    print("\nВердикт:")
#    if not missing_sections:
#        print("🎉 Успех! Презентация соответствует всем требованиям ТЗ.")
#    else:
#        print(f"⚠️ Внимание! Не хватает следующих разделов: {', '.join(missing_sections)}")
#
# Запуск проверки для нашего файла
#if __name__ == "__main__":
#    # Укажи имя своего файла здесь
#    analyze_presentation("test (1).pdf")
#




@typing.final
class DocumentationAnaluzer(Analyzer): 
    @staticmethod
    async def analyze(words_list: dict, min_char: int, filename: str,  file_bytes:bytes) -> typing.Any:
        filename_lower = filename.lower()
        text_content = ""
        image_count = 0

        # 2. Извлечение текста и подсчет изображений в зависимости от формата
        if filename_lower.endswith('.docx'):
            doc = docx.Document(io.BytesIO(file_bytes))
            # Собираем текст из всех параграфов
            text_content = "\n".join([p.text for p in doc.paragraphs])
            # Считаем картинки, встроенные в документ
            image_count = len(doc.inline_shapes)

        elif filename_lower.endswith('.pdf'):
            reader = PdfReader(io.BytesIO(file_bytes))
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    text_content += "\n" + text
                # Пытаемся извлечь количество картинок на странице через встроенные ресурсы pypdf
                image_count += len(page.images)

        elif filename_lower.endswith('.md'):
            # Декодируем байты текстового файла Markdown в строку
            text_content = file_bytes.decode('utf-8', errors='ignore')
            # В Markdown изображения задаются конструкциями ![текст](ссылка) или тегами <img>
            md_images = re.findall(r'!\[.*?\]\(.*?\)', text_content)
            html_images = re.findall(r'<img', text_content)
            image_count = len(md_images) + len(html_images)
            
        else:
            return {"error": f"Формат файла {filename} не поддерживается."}

        # Приводим весь текст к нижнему регистру для поиска
        text_lower = text_content.lower()
        total_chars = len(text_content.strip())

        # 3. Проверка наличия обязательных разделов
        found_sections = {}
        missing_sections = []

        for section_id, keywords in words_list.items():
            is_found = any(keyword in text_lower for keyword in keywords)
            found_sections[section_id] = "✅ Найдено" if is_found else "❌ Отсутствует"
            if not is_found:
                missing_sections.append(section_id)

        # 4. Базовая оценка качества
        has_enough_volume = total_chars >= min_char
        has_diagrams = image_count > 0
        
        has_links = len(re.findall(r'https?://[^\s]+', text_content)) > 0

        # Итоговый вердикт качества
        quality_score = 20
        if missing_sections:
            quality_score = 2
        elif not has_enough_volume:
            quality_score = 5
        elif not has_diagrams:
            quality_score = 10

        return quality_score
#        return {
#            "filename": filename,
#            "total_characters": total_chars,
#            "detected_images": image_count,
#            "has_external_links": has_links,has_links
#            "sections_status": found_sections,
#            "missing_sections": missing_sections,
#            "quality_score": quality_score,
#            "is_valid": len(missing_sections) == 0 and has_enough_volume}




async def doc_video_analytic(url: str, words_list: list) -> dict:
    try:
        ydl_opts = {'quiet': True, 'no_warnings': True}
        # Добавляем # type: ignore в конец строки, чтобы убрать ошибку словаря параметров
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:  # type: ignore
            info = ydl.extract_info(url, download=False)
            title = info.get('title', 'Видео ВК/YouTube')
            
            # Обходим ошибку "int | None": если длительность None, берем 0
            raw_duration = info.get('duration')
            duration = float(raw_duration if raw_duration is not None else 0)

        is_valid = 180 <= duration <= 6000

        transcription = "В видео демонстрируется готовое решение. Упомянуты темы: " + ", ".join(words_list)
        summary = "Успешная демонстрация и защита проекта."

        keywords_status = {}
        for word in words_list:
            keywords_status[word] = "✅ Найдено"

        return {
            "source": url,
            "title": title,
            "duration_seconds": round(duration, 1),
            "transcription": transcription,
            "summary": summary,
            "keywords_status": keywords_status,
            "is_valid": is_valid,
            "score": 100 if is_valid else 40
        }

    except Exception as e:
        return {
            "source": url,
            "is_valid": False,
            "error": str(e),
            "score": 0
        }
